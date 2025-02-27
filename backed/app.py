from flask import Flask, request, jsonify,send_file
import os
import time
import re
import ast
import concurrent.futures
import pymupdf4llm
from flask_cors import CORS
from llama_index.core.llms import ChatMessage
from llama_index.llms.groq import Groq
import groq
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from threading import Lock
import google.generativeai as genai
from flask import send_from_directory
import edge_tts
import asyncio
from pydub import AudioSegment

app = Flask(__name__)
CORS(app)
genai.configure(api_key="AIzaSyBfc6sSOv4dYC7E6cIPsZddWUoQil_wUC0")
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
model = genai.GenerativeModel("gemini-1.5-flash-latest")

api_keys = [
    "Groq_API_KEY1",
    "Groq_API_KEY2",
    "Groq_API_KEY3",
    "Groq_API_KEY4",
    "Groq_API_KEY5"
]
client = groq.Client(api_key="Groq_API_KEY5")

llms = [Groq(model="llama-3.3-70b-versatile", api_key=key) for key in api_keys]

write_lock = Lock()
section_results = {} 

def call_llm_sec(prompt):
    """
    Calls the Groq LLM to get a response for the given prompt.
    """
    response = client.chat.completions.create(
        model="mixtral-8x7b-32768",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_completion_tokens=800,
        top_p=0.9,
        stream=False,
        stop=None,
    )
    return response.choices[0].message.content

def extract_authors(text):
    """
    Extracts authors' names from the text using Groq model.
    If the model fails, it falls back to regex.
    """
    top = text[:1000]
    auth_prompt = f"""Extract the authors' names from the following text and return them strictly as a valid Python list.
Only return the list, nothing else. No extra text.
Text: {top}
[]"""

    response_text = call_llm_sec(auth_prompt)

    try:
        authors = ast.literal_eval(response_text)
        if not isinstance(authors, list):
            raise ValueError("Response is not a valid list.")
    except (SyntaxError, ValueError):
        # Fallback to regex if LLM fails
        authors = re.findall(r"[A-Z][a-z]+(?:\s[A-Z][a-z]+)*", top)
    return authors

def extract_paper_components(text, authors):
    """
    Extracts title, abstract, and sections from markdown text and organizes them into a structured format.
    Ignores subsections like 4.1, 2.3 and only extracts main sections.
    """
    title_pattern = re.compile(r'# (.*?)\n\n')
    abstract_pattern = re.compile(r'Abstract—(.*?)Keywords', re.DOTALL)
    

    section_pattern = re.compile(
        r'(?m)(^\s*(?:[IVXLCDM]+|\d+)(?:\.\s{0,2}|\.\n?|\.?\s{1,2})[A-Z][A-Za-z\- ]*)\n(.*?)(?=\n\s*(?:[IVXLCDM]+|\d+)(?:\.\s{0,2}|\.\n?|\.?\s{1,2})[A-Z]|\n[Rr]eferences|$)', 
        re.DOTALL
    )

    title = title_pattern.search(text).group(1) if title_pattern.search(text) else "Untitled"
    
    TITLE_FILE = os.path.join(OUTPUT_FOLDER, "title.txt")
    ABSTRACT_FILE = os.path.join(OUTPUT_FOLDER, "abstract.txt")
    SLIDES_FILE = os.path.join(OUTPUT_FOLDER, "generated_slides.txt")
    
    # # Remove existing files if they exist
    for file in [TITLE_FILE, ABSTRACT_FILE, SLIDES_FILE]:
        if os.path.exists(file):
            os.remove(file)

    with open(TITLE_FILE, "w", encoding="utf-8") as file:
        file.write(title)

    abstract = abstract_pattern.search(text).group(1).strip() if abstract_pattern.search(text) else ""
    abstract = re.sub(r'\*\*', '', abstract)
    abstract = ' '.join(abstract.split())

    sections = []
    for match in section_pattern.finditer(text):
        section_title = match.group(1).strip()
        section_content = match.group(2).strip()

        # Extract only the main section number (not subsections like 4.1)
        numeral_part = extract_main_section_number(section_title)
        if numeral_part:
            integer_part = numeral_to_int(numeral_part)
            new_section_title = section_title.replace(numeral_part, str(integer_part), 1)
        else:
            integer_part = 999
            new_section_title = section_title
        if(integer_part == 999): continue
        sections.append((integer_part, new_section_title, section_content))
    
    sections.sort(key=lambda x: x[0])

    return {
        'title': title,
        'authors': authors,
        'abstract': abstract,
        'sections': sections
    }

def numeral_to_int(numeral):
    """
    Converts Roman numerals or numeric section numbers to integers.
    """
    roman_dict = {
        'I': 1, 'II': 2, 'III': 3, 'IV': 4, 'V': 5,
        'VI': 6, 'VII': 7, 'VIII': 8, 'IX': 9, 'X': 10
    }
    
    if numeral.isdigit():
        return int(numeral)  # Convert numeric strings like '1', '2' to int
    return roman_dict.get(numeral, 999)  # Convert Roman numerals

def extract_main_section_number(text):
    """
    Extracts only main section numbers (1., 2., I., II.) and ignores subsections like 4.1, 2.3.
    """
    match = re.match(r"^([IVXLCDM]+|\d+)\.", text.strip())  # Ensures only single-level sections
    return match.group(1) if match else None

def call_llm(prompt):
    """
    Calls the Groq LLM and returns structured output.
    """
    response = client.chat.completions.create(
        model="mixtral-8x7b-32768",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_completion_tokens=800,
        top_p=0.9,
    )
    return response.choices[0].message.content

def verify_section_titles(section_titles):
    """
    Sends only section titles to the LLM for verification and correction.
    Ensures that only valid research paper section headers are returned.
    """
    prompt = f"""You are an expert in structuring academic research papers. 
Below is a list of section headers extracted from a research paper. Some of these headers might be incorrect, redundant, or misclassified. 

Your task:
1. Identify and **remove any incorrect or misclassified headers** that do not belong to a research paper.
2. Ensure the remaining headers **follow a logical order** as seen in standard research papers.
3. **Do not modify correct headers**; only remove the invalid ones.
4. Return the cleaned headers as a **valid Python list**. **No explanations, no extra text.** 

Here is the extracted list of section headers:
{section_titles}

**Expected Output Format:**
Example: ["Introduction", "Literature Review", "Methodology", "Results", "Discussion", "Conclusion"]

Your response should be strictly a valid Python list of cleaned section titles."""

    try:
        response_text = call_llm(prompt)
        verified_titles = ast.literal_eval(response_text)
        
        if isinstance(verified_titles, list):
            return verified_titles
    except (SyntaxError, ValueError):
        pass  # If LLM fails, return original titles
    
    return section_titles  # Fallback to original titles

def write_to_txt_file(filename, content):
    """Writes content to a file in a thread-safe manner."""
    with write_lock:
        with open(filename, "a", encoding="utf-8") as file:
            file.write(content + "\n\n")

def main_prompt_small(title, content):
    return f"""You are given a research paper module on the topic *{title}* with the content provided at the end. Your task is to *summarize the content accurately* and *generate structured PowerPoint slides* in a hybrid format (*brief explanation + bullet points*) while maintaining key technical details.

    ## *⚠ Important Rules:*
    - *Strictly include only the summarized research content* – do NOT add any extra text (e.g., "Note:", "Summary:", "Conclusion:").
    - *Remove unnecessary characters* like _, #, @, etc.
    - *Use a hybrid format* – include a brief explanation followed by bullet points.
    - *Optimize content to fit within 1-2 slides per section* (Avoid 3rd slide unless absolutely necessary).
    - *Introduction & Conclusion slides must always be single-page summaries and only in the start and end respectively instead of anywhere in between!!*
    - *Ensure bullet points highlight key concepts, methods, and results concisely.*
    - *Each slide must follow the specified format and contain only research-relevant content.*
    ---
    ## *🔹 Summarization Rules*
    - *Keep it concise and precise* while maintaining all critical technical details.
    - *Ensure numerical data, equations, and key findings are retained.*
    - If the content can fit within *one slide*, do NOT generate an extra slide for minor content.
    - If needed, *spread content across two full slides*, ensuring no incomplete slides.
    - *Use structured bullet points* for clarity while preserving essential details.
    - don't include equation or mathematics formula
    ---
    ## *📌 Slide Generation Instructions:*
    - *Slide 1:* *Title + Short 2-3 line overview of the module* (Must include #Image=True).
    - *Slide 2:* *Main summarized content (explanation + bullet points)* (Must include #Image=False).
    - *Slide 3 (Only if necessary):* *Additional critical details* (Only if Slide 2 exceeds limit, otherwise do NOT generate).
    - *Conclusion Slide:* No image reference required (#Image tag not needed).

    ---
    ## *📝 Expected Response Format (Follow Exactly)*
    ```plaintext
    #Slide: 1
    #Header: {title}
    #Image: True
    #Content:
    [2-3 line concise overview of the module.]

    #Slide: 2
    #Header: {title}
    #Image: False
    #Content:
    [Short explanation of the content, followed by structured bullet points.]
    - Key Point 1
    - Key Point 2
    - Key Point 3 (Include numerical data, equations, or key technical insights)

    ---
    ## *🚀 Here is the Content to Process:*
    {content}
    """

def call_llm_main(prompt):
    """Calls the LLM model and returns generated content."""
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.35,
        max_completion_tokens=1500,
        top_p=0.85,
        stream=False,
        stop=None,
    )
    return response.choices[0].message.content

def process_section(index, section_title, content):
    """Processes each section in parallel and stores the results in order."""
    clean_title = " ".join(section_title.split()[1:]) if section_title[0].isdigit() else section_title
    prompt = main_prompt_small(clean_title, content)

    # Call LLM
    response = call_llm_main(prompt)

    # Store result in dictionary (indexed to preserve order)
    with write_lock:
        section_results[index] = f"#Header: {clean_title}\n{response}"

    print(f"✅ Processed Section: {clean_title}")

def read_file(filename):
    if os.path.exists(filename):
        with open(filename, "r", encoding="utf-8") as file:
            return file.read().strip()
    return None

def parse_llm_output(llm_response):
    slides = re.split(r'#Slide:\s*\d+', llm_response.strip())
    slide_numbers = re.findall(r'#Slide:\s*(\d+)', llm_response)
    parsed_slides = []

    for idx, slide in enumerate(slides[1:]):
        lines = slide.strip().split('\n')
        header_line = next((line for line in lines if line.startswith("#Header:")), None)
        header = header_line.split(": ", 1)[1] if header_line else "No Header"
        
        content_lines = [line for line in lines if line.startswith("#Content:") or not line.startswith("#")]
        content = "\n".join(line.replace("#Content:", "").strip().lstrip('-').strip() for line in content_lines).strip()

        parsed_slides.append({"slide_number": int(slide_numbers[idx]), "title": header, "content": content})

    return parsed_slides

# Function to apply template styles
def apply_template_styles(placeholder, template_choice, is_title=True):
    templates = {
        'simple': {'title_font': 'Calibri Light (Headings)', 'content_font': 'Calibri (Body)',
                   'title_size': 50, 'content_size': 26, 'color': RGBColor(0, 0, 0)},
        'modern': {'title_font': 'Tw Cen MT', 'content_font': 'Segoe UI Light',
                   'title_size': 28, 'content_size': 18, 'color': RGBColor(0, 0, 0)}
    }
    template = templates.get(template_choice, templates['modern'])
    font_name = template['title_font'] if is_title else template['content_font']
    font_size = template['title_size'] if is_title else template['content_size']
    font_color = template['color']

    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color

# Function to create PowerPoint slides
def create_ppt(title, abstract, slides_content, template_choice):
    prs = Presentation()

    # Slide 1: Title Slide
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        apply_template_styles(title_placeholder, template_choice, is_title=True)

    if abstract:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Abstract"
        apply_template_styles(slide.shapes.title, template_choice, is_title=True)
        slide.shapes.placeholders[1].text = abstract
        apply_template_styles(slide.shapes.placeholders[1], template_choice, is_title=False)

    for slide_content in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_content['title']
        apply_template_styles(slide.shapes.title, template_choice, is_title=True)
        slide.shapes.placeholders[1].text = slide_content['content']
        apply_template_styles(slide.shapes.placeholders[1], template_choice, is_title=False)

    output_path = os.path.join(OUTPUT_FOLDER, "generated_presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved as {output_path}")
    return output_path




def extract_text_from_pdf(pdf_path, output_path):
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"Error: File '{pdf_path}' not found.")

    md_text = pymupdf4llm.to_markdown(os.path.abspath(pdf_path))

    with open(output_path, "w", encoding="utf-8") as file:
        file.write(md_text)

def split_text(text, max_tokens=8000):
    return [" ".join(m.group().split()) for m in re.finditer(r'(\S+\s*){1,%d}' % max_tokens, text)]

def gemini_summ(text, target_sentences=40, retries=2, single_paragraph=False):
    prompt = (
        f"Summarize the following research text in a single paragraph containing approximately {target_sentences} sentences. "
        "Ensure the paragraph is well-structured, cohesive, and captures all critical insights, methodologies, and conclusions "
        "without listing points separately. The summary should be fluid and connected, avoiding bullet points or line breaks."
    )

    for _ in range(retries):
        try:
            response = model.generate_content([prompt, text])
            return response.text.strip()
        except Exception as e:
            print(f"Retrying due to error: {e}")
            time.sleep(3)  # Shorter wait time

    return "Summarization failed."

# def process_paragraph(idx, paragraph):
#     llm_index = idx % len(llms)
#     return idx, llama_summ(llms[llm_index], paragraph, target_lines=2)

def summarize_text(input_path, summary_path):
    start_time = time.time()
    
    with open(input_path, "r", encoding="utf-8") as file:
        md_text = file.read()

    paragraphs = split_text(md_text)
    results = {}

    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as executor:  # Parallel execution
        future_to_index = {executor.submit(gemini_summ, para, 10): idx for idx, para in enumerate(paragraphs)}

        for future in concurrent.futures.as_completed(future_to_index):
            idx = future_to_index[future]
            try:
                results[idx] = future.result()
            except Exception:
                results[idx] = "Failed to summarize."

    first_pass_summary = " ".join([results[i] for i in sorted(results)])

    final_summary = gemini_summ(first_pass_summary, target_sentences=40, single_paragraph=True)

    with open(summary_path, "w", encoding="utf-8") as file:
        file.write(final_summary)

    end_time = time.time()
    print(f"Summarization completed in {end_time - start_time:.2f} seconds.")

    return final_summary

async def generate_speech(text, voice, filename):
    tts = edge_tts.Communicate(text=text, voice=voice)
    await tts.save(filename)
    print(f"Generated: {filename}")  # Debugging

def merge_audio(audio_files, intro_music, outro_music, output_filename):
    if not audio_files:
        print("❌ No audio files generated. Skipping merge.")
        return None

    segments = []
    for file in audio_files:
        try:
            segment = AudioSegment.from_mp3(file)
            segments.append(segment)
        except Exception as e:
            print(f"❌ Error loading {file}: {e}")

    if not segments:
        print("❌ No valid audio segments to merge.")
        return None

    # Merge all segments
    try:
        podcast = sum(segments)
    except Exception as e:
        print(f"❌ Error merging audio segments: {e}")
        return None

    # Load intro and outro music
    try:
        intro = AudioSegment.from_file(intro_music) - 10
        outro = AudioSegment.from_file(outro_music) - 10
    except Exception as e:
        print(f"❌ Failed to load intro/outro music: {e}")
        return None

    # Final assembly
    try:
        final_podcast = intro + podcast + outro
    except Exception as e:
        print(f"❌ Error assembling final podcast: {e}")
        return None

    # Export as MP3
    try:
        final_podcast.export(output_filename, format="mp3")
        print(f"✅ Podcast saved at: {output_filename}")
    except Exception as e:
        print(f"❌ Error exporting final podcast: {e}")
        return None

    # Cleanup temporary files
    for file in audio_files:
        try:
            os.remove(file)
        except Exception as e:
            print(f"❌ Error deleting {file}: {e}")

    return output_filename




async def process_script(script_lines):
    
    # Define voices
    voices = {"Host": "en-GB-SoniaNeural", "Guest": "en-US-JessaNeural"}

    # List to store generated file names
    audio_files = []
    tasks = []

    for i, line in enumerate(script_lines):
        line = line.strip()
        if not line:
            continue  # Skip empty lines

        for role in voices.keys():
            if line.startswith(f"{role}:"):
                text = line.replace(f"{role}:", "").strip()
                filename = f"{role.lower()}_{i}.mp3"
                audio_files.append(filename)
                tasks.append(generate_speech(text, voices[role], filename))
                break  # Avoid processing line multiple times

    if tasks:
        await asyncio.gather(*tasks)

    return audio_files




@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    pdf_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pdf_path)

    output_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}.txt")

    try:
        extract_text_from_pdf(pdf_path, output_file_path)

        with open(output_file_path, "r", encoding="utf-8") as f:
            extracted_text = f.read()
        
        return jsonify({
            "message": "Text extracted successfully",
            "text": extracted_text,
            "filename": file.filename
        }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/summarize', methods=['POST'])
def summarize():
    data = request.json
    filename = data.get("filename")

    if not filename:
        return jsonify({"error": "Filename not provided"}), 400

    input_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(filename)[0]}.txt")
    summary_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(filename)[0]}_summary.txt")

    if not os.path.exists(input_path):
        return jsonify({"error": "Extracted text file not found"}), 404

    try:
        summary_text = summarize_text(input_path, summary_path)
        print(summary_text)
        return jsonify({
            "message": "Summarization completed",
            "summary": summary_text
        }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    # Get file name from frontend request
    data = request.get_json()
    file_name = data.get("filename")
    print(file_name)
    if not file_name:
        return jsonify({"error": "File name not provided"}), 400

    # Ensure the file path exists in OUTPUT_FOLDER
    file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file_name)[0]}.txt")

    print(file_path)
    if not os.path.exists(file_path):
        return jsonify({"error": f"File '{file_name}.txt' not found"}), 404

    # Read text from the uploaded file
    with open(file_path, "r", encoding="utf-8") as f:
        md_text = f.read()
    print(md_text[:100])
    if not md_text.strip():
        return jsonify({"error": "Extracted text is empty"}), 400

    # Extract authors and paper components
    authors = extract_authors(md_text)
    paper_components = extract_paper_components(md_text, authors)

    section_titles = [section_title for _, section_title, content in paper_components['sections']]
    print(section_titles)
    verified_titles = verify_section_titles(section_titles)
    print(f"verified titles {verified_titles}")
    paper_components["sections"] = [
    (num, verified_title, content)
    for num, section_title, content in paper_components["sections"]
    for verified_title in verified_titles
    if verified_title.strip().lower() in section_title.strip().lower()
]



    print("\nFinal Verified Sections:")
    for _, section_title, content in paper_components['sections']:
        print(f"Section: {section_title}")

    # Prepare output file
    output_file = os.path.join(OUTPUT_FOLDER, "generated_slides.txt")
    if os.path.exists(output_file):
        os.remove(output_file)
    print("hello")
    # Process sections in parallel
    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as executor:  # Increase to 8 workers for speed
        futures = [executor.submit(process_section, idx, title, content) for idx, (_, title, content) in enumerate(paper_components['sections'])]

        # Wait for all threads to complete
        concurrent.futures.wait(futures)

# Write results in order after processing
    with open(output_file, "w", encoding="utf-8") as file:
        for index in sorted(section_results.keys()):
            file.write(section_results[index] + "\n\n")

    print("🚀✅ All sections processed and written in order successfully!")
    # Define additional files in OUTPUT_FOLDER
    TITLE_FILE = os.path.join(OUTPUT_FOLDER, "title.txt")
    ABSTRACT_FILE = os.path.join(OUTPUT_FOLDER, "abstract.txt")
    SLIDES_FILE = os.path.join(OUTPUT_FOLDER, "generated_slides.txt")
    
    

    
    template_choice = "modern"

    # Read text files
    title_text = read_file(TITLE_FILE) if os.path.exists(TITLE_FILE) else ""
    abstract_text = read_file(ABSTRACT_FILE) if os.path.exists(ABSTRACT_FILE) else ""
    slides_text = read_file(SLIDES_FILE) if os.path.exists(SLIDES_FILE) else ""
    print(title_text)
    slides_content = parse_llm_output(slides_text) if slides_text else []

    # Generate PPT
    ppt_path = create_ppt(title_text, abstract_text, slides_content, template_choice)
    print(ppt_path)
    ppt_filename = os.path.basename(ppt_path)
    print(ppt_filename)
    return jsonify({
        "message": "✅ Presentation generated successfully!",
        "ppt_path": ppt_filename
    }), 200

@app.route("/generate_podcast", methods=["POST"])
def generate_podcast():
    start_time = time.time()
    data = request.json
    filename = data.get("filename")

    if not filename:
        return jsonify({"error": "Filename not provided"}), 400
    print(filename)

    input_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(filename)[0]}_summary.txt")
    podcast_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(filename)[0]}_podcast.mp3")


    try:
            with open(input_path, "r", encoding="utf-8") as f:
                final_summary = f.read()
    except FileNotFoundError:
            print(f"❌ Summary file not found: {input_path}")
            return jsonify({"error": "Summary file not found"}), 404

    prompt = """You are an AI that converts research paper summaries into a podcast script.
    The podcast has two speakers:
    1. Host: Asks engaging, curious questions in more detail.
    2. Guest: Explains concepts in a casual, educational manner.

    The tone should be friendly, engaging, and easy to understand.
    The podcast should start with an introduction where the host welcomes listeners and introduces the topic.
    Then, the host asks insightful questions, and the guest responds.
    Finally, it should end with a friendly outro encouraging curiosity.

    Podcast Format:
    Host: Welcome to [Podcast Name]! Today, we’re discussing [Topic].
    Guest: [Introduction to the topic]
    Host: [First Question]
    Guest: [Answer]
    ...
    Host: Thank you for joining us today! See you in the next episode.

    Do not include unnecessary text. The output should contain only 
    Host: question 
    guest: answer format  without special characters like * or #.

    Now, generate a script from the following research paper summary:
    """

    # Combine the prompt with the summary
    full_prompt = prompt + "\n\n" + final_summary

    # Generate podcast script
    
    try:
        response = call_llm_main(full_prompt)
        podcast_script = response.strip()

        if not podcast_script:
            raise Exception("Generated podcast script is empty.")
    except Exception as e:
        print(f"❌ Podcast script generation failed: {e}")
        
    final_path = None
    intro_music = os.path.join(OUTPUT_FOLDER, "Rome.mp3")
    outro_music = os.path.join(OUTPUT_FOLDER, "Barcelona.mp3")
    try:
        script_lines = podcast_script.split("\n")
        print(script_lines)
    # Ensure we are running in a new event loop
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        audio_files = loop.run_until_complete(process_script(script_lines))

        if not audio_files:
            print("❌ No audio files were generated.")
            # return jsonify({"error": "Podcast script processing failed"}), 500
        final_path = merge_audio(audio_files, intro_music, outro_music, podcast_path)
    except Exception as e:
        print(f"❌ Podcast audio generation failed: ")
        # return jsonify({"error": f"Podcast audio generation failed: {e}"}), 500

    end_time = time.time()
    execution_time = round(end_time - start_time, 2)

    print(f"✅ Podcast generated: {final_path}")
    print(f"⏳ Execution Time: {execution_time} seconds")

    return jsonify({
        "message": "Podcast generated successfully",
        "podcast_path": f"{OUTPUT_FOLDER}/{os.path.splitext(filename)[0]}_podcast.mp3",
        "execution_time": execution_time
    })

@app.route("/output/<filename>")
def download_podcast(filename):
    """Serve the podcast file"""
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)


@app.route("/download/<filename>", methods=["GET"])
def download_ppt(filename):
    file_path = os.path.join(OUTPUT_FOLDER, filename)

    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404

    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)


if __name__ == '__main__':
    app.run(debug=True)
